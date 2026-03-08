"""
📄 PDF → WORD PPTX Converter
"""

import streamlit as st
import io
import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import Image, ImageDraw
from collections import Counter
from docx import Document
from docx.shared import Inches, Pt as DocxPt, RGBColor as DocxRGBColor

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PDF → DOCX - PPTX",
    page_icon="📄",
    layout="centered",
)

# ── Styling ───────────────────────────────────────────────────────────────────
st.markdown("""
<style>
/* Upload area */
[data-testid="stFileUploader"] {
    border: 2px dashed #4A90D9;
    border-radius: 12px;
    padding: 1rem;
}
/* Settings card */
.settings-card {
    background: #1e1e2e;
    border-radius: 12px;
    padding: 1.2rem 1.5rem;
    margin-bottom: 1rem;
}
/* Status badge */
.badge {
    display: inline-block;
    padding: 2px 10px;
    border-radius: 20px;
    font-size: 0.78rem;
    font-weight: 600;
}
.badge-native { background:#2ecc71; color:#000; }
.badge-ocr    { background:#e67e22; color:#fff; }
</style>
""", unsafe_allow_html=True)

# ── Header ────────────────────────────────────────────────────────────────────
st.title("📄 PDF Converter")
st.caption("Converte ogni pagina PDF in PPTX o DOCX con testo selezionabile. ⚠️ Non esegue OCR: funziona solo con PDF che contengono testo nativo.")
st.image('aa.jpg', use_container_width=True)
st.divider()

# ── File upload ───────────────────────────────────────────────────────────────
uploaded = st.file_uploader(
    "Trascina qui il tuo PDF oppure clicca per sfogliare",
    type=["pdf"],
    accept_multiple_files=False,
    label_visibility="visible",
    key="pdf_upload",
)

# ── Settings ──────────────────────────────────────────────────────────────────
st.subheader("⚙️ Impostazioni")

output_format = st.radio(
    "Formato di output",
    options=["PPTX", "DOCX"],
    index=0,
    horizontal=True,
)

col1, col2 = st.columns(2)

with col1:
    if output_format == "PPTX":
        dpi = st.select_slider(
            "Qualità immagine (DPI)",
            options=[72, 100, 150, 200, 250, 300],
            value=150,
            help="DPI per le pagine scansionate (senza testo nativo).",
        )
    else:
        dpi = 150
        st.empty()

with col2:
    ext = ".pptx" if output_format == "PPTX" else ".docx"
    default_name = f"file_convertito{ext}"
    output_name = st.text_input(
        "Nome file di output",
        value=default_name,
        help="Il file sarà scaricabile con questo nome.",
    )
    if not output_name.endswith(ext):
        output_name = output_name.rsplit(".", 1)[0] + ext

st.divider()

# ── Conversion logic ──────────────────────────────────────────────────────────

def convert(pdf_bytes: bytes, dpi: int) -> bytes:
    """Core conversion — runs inside Streamlit, returns PPTX bytes."""


    EMU_PER_INCH = 914_400
    EMU_PER_PT   = EMU_PER_INCH / 72

    def pts_to_emu(v):  return int(v * EMU_PER_PT)

    def int_color_to_rgb(c):
        """Convert fitz integer color (0xRRGGBB) to RGBColor."""
        r = (c >> 16) & 0xFF
        g = (c >> 8)  & 0xFF
        b =  c        & 0xFF
        return RGBColor(r, g, b)

    def clean_font_name(raw):
        """Strip style suffixes embedded in PDF font names."""
        name = raw.split("+")[-1]          # remove subset prefix (e.g. ABCDEF+Arial)
        name = name.replace("-", " ").split(",")[0].strip()
        return name or "Calibri"

    def render_bg_no_text(page):
        """Render page with text erased using the sampled background colour per word."""


        scale  = dpi / 72
        pix    = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        img    = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        draw   = ImageDraw.Draw(img)
        W, H   = img.size
        bdr    = max(2, int(scale))   # border thickness in px (≥2 px)

        for word in page.get_text("words"):
            x0, y0, x1, y1 = word[:4]
            px0, py0 = int(x0 * scale), int(y0 * scale)
            px1, py1 = int(x1 * scale), int(y1 * scale)

            # Collect pixels in a thin strip OUTSIDE the word bbox.
            # These are background pixels — text glyphs are inside the bbox.
            samples: list = []
            if py0 - bdr >= 0:
                samples += list(img.crop((max(0,px0), py0-bdr, min(W,px1), py0)).getdata())
            if py1 + bdr <= H:
                samples += list(img.crop((max(0,px0), py1, min(W,px1), min(H,py1+bdr))).getdata())
            if px0 - bdr >= 0:
                samples += list(img.crop((px0-bdr, max(0,py0), px0, min(H,py1))).getdata())
            if px1 + bdr <= W:
                samples += list(img.crop((px1, max(0,py0), min(W,px1+bdr), min(H,py1))).getdata())

            # Most-common pixel colour = dominant background (not the glyph colour)
            fill = Counter(samples).most_common(1)[0][0] if samples else (255, 255, 255)

            draw.rectangle(
                [max(0, px0-1), max(0, py0-1), min(W-1, px1+1), min(H-1, py1+1)],
                fill=fill,
            )

        buf = io.BytesIO()
        img.save(buf, "PNG")
        buf.seek(0)
        return buf

    def add_native_page(slide, page, emu_w, emu_h):
        """Populate a slide with native text spans and embedded images from page."""
        blocks = page.get_text(
            "dict",
            flags=fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_MEDIABOX_CLIP,
        )["blocks"]
        word_count = 0

        for block in blocks:
            btype = block.get("type", 0)

            # ── Embedded image ────────────────────────────────────────────────
            if btype == 1:
                img_bytes = block.get("image")
                if not img_bytes:
                    continue
                x0, y0, x1, y1 = block["bbox"]
                left   = max(0, pts_to_emu(x0))
                top    = max(0, pts_to_emu(y0))
                width  = max(1, min(pts_to_emu(x1 - x0), emu_w - left))
                height = max(1, min(pts_to_emu(y1 - y0), emu_h - top))
                try:
                    slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
                except Exception:
                    pass
                continue

            # ── Text block ────────────────────────────────────────────────────
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    txt = span.get("text", "")
                    if not txt.strip():
                        continue

                    x0, y0, x1, y1 = span["bbox"]
                    left   = max(0, pts_to_emu(x0))
                    top    = max(0, pts_to_emu(y0))
                    width  = max(1, min(pts_to_emu(x1 - x0), emu_w - left))
                    height = max(1, min(pts_to_emu(y1 - y0), emu_h - top))

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = False
                    tf.auto_size = None

                    run = tf.paragraphs[0].add_run()
                    run.text = txt
                    word_count += len(txt.split())

                    font_size = span.get("size", 12)
                    run.font.size = Pt(font_size)

                    flags = span.get("flags", 0)
                    run.font.bold   = bool(flags & 16)
                    run.font.italic = bool(flags & 2)

                    color = span.get("color", 0)
                    run.font.color.rgb = int_color_to_rgb(color)

                    font_name = span.get("font", "")
                    if font_name:
                        run.font.name = clean_font_name(font_name)

        return word_count

    def add_ocr_page(slide, page, emu_w, emu_h):
        """Fallback for scanned pages: place the page image, no text overlay."""
        matrix   = fitz.Matrix(dpi / 72, dpi / 72)
        pix      = page.get_pixmap(matrix=matrix, alpha=False)
        img_io   = io.BytesIO(pix.tobytes("png"))
        slide.shapes.add_picture(img_io, 0, 0, emu_w, emu_h)
        return 0  # no text words added

    # ── Open PDF ──────────────────────────────────────────────────────────────
    doc   = fitz.open(stream=pdf_bytes, filetype="pdf")
    n     = len(doc)
    first = doc[0]

    prs = Presentation()
    prs.slide_width  = pts_to_emu(first.rect.width)
    prs.slide_height = pts_to_emu(first.rect.height)
    blank = prs.slide_layouts[6]

    progress  = st.progress(0, text="Conversione in corso…")
    info_rows = []

    for i in range(n):
        page  = doc[i]
        emu_w = pts_to_emu(page.rect.width)
        emu_h = pts_to_emu(page.rect.height)

        slide = prs.slides.add_slide(blank)

        raw     = page.get_text("text").strip()
        is_native = len(raw) >= 10

        if is_native:
            bg = render_bg_no_text(page)
            slide.shapes.add_picture(bg, 0, 0, emu_w, emu_h)
            words = add_native_page(slide, page, emu_w, emu_h)
            mode  = "native"
        else:
            words = add_ocr_page(slide, page, emu_w, emu_h)
            mode  = "OCR (immagine)"

        info_rows.append((i + 1, mode, words))
        progress.progress((i + 1) / n, text=f"Pagina {i+1}/{n}  [{mode}]")

    progress.empty()
    doc.close()

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), info_rows


# ── Conversion logic — DOCX ───────────────────────────────────────────────────

def convert_to_docx(pdf_bytes: bytes, dpi: int):
    """Converte PDF in DOCX — testo nativo o immagine per pagine scansionate."""

    def clean_font_name(raw):
        name = raw.split("+")[-1]
        name = name.replace("-", " ").split(",")[0].strip()
        return name or "Calibri"

    def int_to_docx_rgb(c):
        r = (c >> 16) & 0xFF
        g = (c >> 8)  & 0xFF
        b =  c        & 0xFF
        return DocxRGBColor(r, g, b)

    doc      = fitz.open(stream=pdf_bytes, filetype="pdf")
    n        = len(doc)
    word_doc = Document()

    for section in word_doc.sections:
        section.top_margin    = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin   = Inches(0.75)
        section.right_margin  = Inches(0.75)

    progress  = st.progress(0, text="Conversione in corso…")
    info_rows = []

    for i in range(n):
        page       = doc[i]
        raw        = page.get_text("text").strip()
        is_native  = len(raw) >= 10
        word_count = 0

        if is_native:
            blocks = page.get_text(
                "dict",
                flags=fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_MEDIABOX_CLIP,
            )["blocks"]

            for block in blocks:
                btype = block.get("type", 0)

                if btype == 1:
                    img_bytes = block.get("image")
                    if img_bytes:
                        try:
                            para = word_doc.add_paragraph()
                            run  = para.add_run()
                            run.add_picture(io.BytesIO(img_bytes), width=Inches(5))
                        except Exception:
                            pass
                    continue

                for line in block.get("lines", []):
                    para = word_doc.add_paragraph()
                    for span in line.get("spans", []):
                        txt = span.get("text", "")
                        if not txt:
                            continue
                        run = para.add_run(txt)
                        word_count += len(txt.split())

                        font_size = span.get("size", 12)
                        run.font.size = DocxPt(font_size)

                        flags = span.get("flags", 0)
                        run.font.bold   = bool(flags & 16)
                        run.font.italic = bool(flags & 2)

                        color = span.get("color", 0)
                        run.font.color.rgb = int_to_docx_rgb(color)

                        font_name = span.get("font", "")
                        if font_name:
                            run.font.name = clean_font_name(font_name)

            mode = "native"
        else:
            pix    = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), alpha=False)
            img_io = io.BytesIO(pix.tobytes("png"))
            para   = word_doc.add_paragraph()
            run    = para.add_run()
            run.add_picture(img_io, width=Inches(6))
            mode = "OCR (immagine)"

        info_rows.append((i + 1, mode, word_count))
        progress.progress((i + 1) / n, text=f"Pagina {i+1}/{n}  [{mode}]")

        if i < n - 1:
            word_doc.add_page_break()

    progress.empty()
    doc.close()

    out = io.BytesIO()
    word_doc.save(out)
    return out.getvalue(), info_rows


# ── Convert button ────────────────────────────────────────────────────────────
if uploaded is not None:
    st.success(f"File caricato: **{uploaded.name}** ({uploaded.size / 1024:.1f} KB)")

    if st.button(f"🚀 Converti in {output_format}", type="primary", use_container_width=True):
        try:
            pdf_bytes = uploaded.read()
            with st.spinner("Elaborazione in corso…"):
                if output_format == "PPTX":
                    result_bytes, info = convert(pdf_bytes, dpi)
                    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                else:
                    result_bytes, info = convert_to_docx(pdf_bytes, dpi)
                    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

            st.success("✅ Conversione completata!")

            # Page summary table
            st.subheader("📋 Riepilogo pagine")
            col_a, col_b, col_c = st.columns(3)
            col_a.metric("Pagine totali", len(info))
            col_b.metric("Con testo nativo", sum(1 for _, m, _ in info if m == "native"))
            col_c.metric("Con OCR (immagine)", sum(1 for _, m, _ in info if m != "native"))

            with st.expander("Dettaglio per pagina"):
                rows_html = "".join(
                    f"<tr><td style='padding:4px 12px;'>Pagina {p}</td>"
                    f"<td><span class='badge {'badge-native' if m == 'native' else 'badge-ocr'}'>{m.upper()}</span></td>"
                    f"<td style='padding:4px 12px;'>{w} parole</td></tr>"
                    for p, m, w in info
                )
                st.markdown(
                    f"<table style='width:100%;border-collapse:collapse;'>{rows_html}</table>",
                    unsafe_allow_html=True,
                )

            # Download
            st.download_button(
                label=f"⬇️ Scarica {output_format}",
                data=result_bytes,
                file_name=output_name,
                mime=mime,
                use_container_width=True,
                type="primary",
            )

            # Reset
            if st.button("🗑️ Nuova conversione", use_container_width=True):
                del st.session_state["pdf_upload"]
                st.rerun()

        except Exception as e:
            st.error(f"Errore durante la conversione:\n\n```\n{e}\n```")

else:
    st.info("⬆️ Carica un PDF per iniziare.")

# ── Footer ────────────────────────────────────────────────────────────────────
st.divider()
st.caption(
    "💡 **Suggerimento:** DPI 150 è un buon compromesso qualità/velocità. "
    "Le pagine senza testo nativo (scansioni) vengono inserite come immagine."
)
